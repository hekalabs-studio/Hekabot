import sys
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
# Kalau ada folder bin/poppler/ di project ini, pakai itu duluan (gak perlu edit PATH sistem sama
# sekali -- sama kayak pola ffmpeg/yt-dlp/Real-ESRGAN di bot ini). Kalau gak ada, fallback ke PATH.
_POPPLER_BIN = os.path.join(PROJECT_ROOT, "bin", "poppler")
POPPLER_PATH = _POPPLER_BIN if os.path.isdir(_POPPLER_BIN) else None

def main():
    input_path, output_path = sys.argv[1], sys.argv[2]
    images = convert_from_path(input_path, dpi=150, poppler_path=POPPLER_PATH)
    if not images:
        raise RuntimeError("PDF tidak punya halaman untuk dikonversi.")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(10 * images[0].height / images[0].width)
    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmp_dir:
        for i, img in enumerate(images):
            slide = prs.slides.add_slide(blank_layout)
            img_path = os.path.join(tmp_dir, f"page_{i}.png")
            img.save(img_path)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(output_path)

if __name__ == "__main__":
    main()
